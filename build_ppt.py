"""Build the Task 5 PPT deck (max 10 slides) from analysis results."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import json
from pathlib import Path

ROOT = Path(__file__).parent
FIG = ROOT / "reports" / "figures"
MODELS = ROOT / "models"

with open(MODELS / "logreg_metrics.json") as f:
    lr = json.load(f)
with open(MODELS / "lightgbm_metrics.json") as f:
    lgb = json.load(f)

# ---- validated palette slots (dataviz skill reference palette, light mode) ----
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
RED = RGBColor(0xE3, 0x49, 0x48)      # slot 8 - accent / Home Credit thematic red
BLUE = RGBColor(0x2A, 0x78, 0xD6)     # slot 1 - Logistic Regression
ORANGE = RGBColor(0xEB, 0x68, 0x34)   # slot 2 - LightGBM
AQUA = RGBColor(0x1B, 0xAF, 0x7A)     # slot 3 - "good" status
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEE, 0xED, 0xEA)

FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=SURFACE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def textbox(slide, left, top, width, height, text, size=14, bold=False, color=INK,
            align=PP_ALIGN.LEFT, font=FONT, italic=False, line_spacing=1.0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return tb


def header(slide, kicker, title):
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    textbox(slide, Inches(0.55), Inches(0.32), Inches(11.5), Inches(0.35), kicker,
            size=13, bold=True, color=RED)
    textbox(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.7), title,
            size=28, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.28), Inches(12.2), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()


def footer(slide, page_num):
    textbox(slide, Inches(0.55), Inches(7.12), Inches(6), Inches(0.3),
            "Home Credit Default Risk — Rakamin x Home Credit Indonesia Virtual Internship, Task 5",
            size=9, color=INK_SECONDARY)
    textbox(slide, Inches(12.4), Inches(7.12), Inches(0.6), Inches(0.3), str(page_num),
            size=9, color=INK_SECONDARY, align=PP_ALIGN.RIGHT)


def bullets(slide, left, top, width, height, items, size=15, color=INK, space_after=10, bold_lead=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = "▪ " + lead
            r1.font.bold = True; r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest if rest.startswith(" ") else " " + rest
            r2.font.size = Pt(size); r2.font.color.rgb = INK_SECONDARY; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = "▪ " + item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def stat_tile(slide, left, top, width, height, value, label, accent):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LIGHT_GRAY; box.line.width = Pt(1)
    box.adjustments[0] = 0.06
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = value
    r0.font.size = Pt(34); r0.font.bold = True; r0.font.color.rgb = accent; r0.font.name = FONT
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = label
    r1.font.size = Pt(12); r1.font.color.rgb = INK_SECONDARY; r1.font.name = FONT
    return box


def picture_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    l = left + int((max_w - w) / 2)
    t = top + int((max_h - h) / 2)
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)


# ============================================================ SLIDE 1 — TITLE
s = add_slide(); set_bg(s)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.35))
bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5), "HOME CREDIT DEFAULT RISK",
        size=18, bold=True, color=RED)
textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.7),
        "From Risk Prediction\nto Business Action", size=40, bold=True, color=INK, line_spacing=1.05)
textbox(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.6),
        "Approving the right customers, on the right terms, more often",
        size=18, color=INK_SECONDARY, italic=True)
textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
        "Rakamin Academy x Home Credit Indonesia — Virtual Internship, Task 5",
        size=13, color=INK_SECONDARY)
textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
        "Steven Glecy Haloho", size=12, bold=True, color=INK)

# ============================================================ SLIDE 2 — PROBLEM
s = add_slide(); set_bg(s)
header(s, "01 — THE PROBLEM", "A scoring model has to do two jobs at once")
bullets(s, Inches(0.55), Inches(1.6), Inches(11.8), Inches(2.4), [
    ("Don't reject good customers.", "Applicants capable of repaying should not be turned away — every"
     " false rejection is lost revenue and a customer Home Credit could have served well."),
    ("Don't approve loans set up to fail.", "Applicants at genuine risk of default need to be identified"
     " early — but the goal is prevention, not just exclusion."),
    ("Structure the loan to help the customer succeed.", "Principal, tenor, and repayment calendar should"
     " match the customer's real risk and cash-flow profile — not a one-size-fits-all product."),
], size=17, space_after=18)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.15))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Objective: "
r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = RED; r.font.name = FONT
r2 = p.add_run()
r2.text = ("build ≥2 ML models (incl. Logistic Regression) that rank-order default risk well enough to "
           "support both accept/reject decisions and risk-based loan structuring — evaluated on ROC-AUC / "
           "KS at the model level, and on defaults-avoided vs. good-payers-affected at the business level.")
r2.font.size = Pt(15); r2.font.color.rgb = INK; r2.font.name = FONT
footer(s, 2)

# ============================================================ SLIDE 3 — DATASET
s = add_slide(); set_bg(s)
header(s, "02 — THE DATASET", "One loan record, five behavioral histories")
textbox(s, Inches(0.55), Inches(1.55), Inches(5.6), Inches(0.4), "application_train.csv (main table)",
        size=15, bold=True, color=INK)
bullets(s, Inches(0.55), Inches(2.0), Inches(5.7), Inches(1.6), [
    "307,511 loans, 122 raw columns",
    "TARGET: 1 = defaulted (8.07%), 0 = repaid (91.93%)",
    "Heavily imbalanced — accuracy is not a usable metric",
], size=14, space_after=10)
stat_tile(s, Inches(0.55), Inches(3.75), Inches(2.7), Inches(1.3), "307,511", "loan applications", BLUE)
stat_tile(s, Inches(3.4), Inches(3.75), Inches(2.7), Inches(1.3), "8.07%", "portfolio default rate", RED)

textbox(s, Inches(6.55), Inches(1.55), Inches(6.2), Inches(0.4), "5 supporting behavioral tables",
        size=15, bold=True, color=INK)
bullets(s, Inches(6.55), Inches(2.0), Inches(6.2), Inches(3.3), [
    ("bureau.csv + bureau_balance.csv", " — credit history at other institutions (1.7M / 27.3M rows)"),
    ("previous_application.csv", " — client's past applications to Home Credit (1.67M rows)"),
    ("POS_CASH_balance.csv", " — past point-of-sale / cash loan behavior (10.0M rows)"),
    ("credit_card_balance.csv", " — past Home Credit credit card behavior (3.8M rows)"),
    ("installments_payments.csv", " — installment-level repayment history (13.6M rows)"),
], size=13.5, space_after=12)
textbox(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.3),
        "Engineering challenge: ~2.3GB raw, 8 supporting tables → aggregated to one row per client "
        "(bureau/previous-application history → SK_ID_CURR) and merged with engineered ratios "
        "(credit/income, annuity/income, EXT_SOURCE combinations) into a 307,511 × 196 master feature set.",
        size=13, color=INK_SECONDARY, italic=True)
footer(s, 3)

# ============================================================ SLIDE 4 — INSIGHT 1
s = add_slide(); set_bg(s)
header(s, "03 — TOP INSIGHT #1", "Low-risk segments are under-represented in the portfolio")
picture_fit(s, FIG / "insight_income_type.png", Inches(0.55), Inches(1.55), Inches(6.4), Inches(4.6))
tb = textbox(s, Inches(7.2), Inches(1.6), Inches(5.6), Inches(0.3), "What the data shows", size=14, bold=True, color=RED)
bullets(s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(2.6), [
    "State servants: 5.76% default (vs 8.07% avg) — only 7.1% of applicants",
    "Pensioners: 5.39% default, the lowest of all segments — only 18.0% of applicants",
    "Working (private sector): 9.59% default — yet 51.6% of the whole portfolio",
], size=14, space_after=10)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.5), Inches(5.6), Inches(2.0))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Action: "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = AQUA; r.font.name = FONT
r2 = p.add_run()
r2.text = ("Run targeted acquisition (partnerships with government/pension-disbursing institutions) toward "
           "State servants & Pensioners. This shifts portfolio mix toward lower risk without tightening "
           "underwriting on anyone — a volume-and-quality win, not just a rejection rule.")
r2.font.size = Pt(14); r2.font.color.rgb = INK; r2.font.name = FONT
footer(s, 4)

# ============================================================ SLIDE 5 — INSIGHT 2
s = add_slide(); set_bg(s)
header(s, "04 — TOP INSIGHT #2", "Bureau score should shape loan terms, not just approval")
picture_fit(s, FIG / "insight_ext_source.png", Inches(0.55), Inches(1.55), Inches(6.4), Inches(4.6))
tb = textbox(s, Inches(7.2), Inches(1.6), Inches(5.6), Inches(0.3), "What the data shows", size=14, bold=True, color=RED)
bullets(s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(1.6), [
    "EXT_SOURCE_MEAN bottom quartile: 17.3% default",
    "EXT_SOURCE_MEAN top quartile: 2.7% default — a 6.4x spread",
    "Confirmed as the single strongest driver in both models",
], size=14, space_after=10)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(3.75), Inches(5.6), Inches(2.75))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Action: "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = AQUA; r.font.name = FONT
r2 = p.add_run()
r2.text = ("Don't use this for a binary cutoff alone. Use the full model score to set risk-tiered principal, "
           "tenor, and repayment calendar: top-tier clients get a higher principal ceiling and longer tenor "
           "(volume Home Credit is currently leaving on the table); higher-risk clients get a smaller "
           "principal, shorter tenor, and more frequent (weekly/bi-weekly) instalments — easier to sustain "
           "and a path to a bigger loan next cycle.")
r2.font.size = Pt(13.5); r2.font.color.rgb = INK; r2.font.name = FONT
footer(s, 5)

# ============================================================ SLIDE 6 — MODEL APPROACH
s = add_slide(); set_bg(s)
header(s, "05 — MODELING APPROACH", "Final models: what features, what preprocessing, what algorithm")

col_w = Inches(5.9)
# Logistic Regression card
box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.55), col_w, Inches(5.1))
box1.fill.solid(); box1.fill.fore_color.rgb = WHITE; box1.line.color.rgb = BLUE; box1.line.width = Pt(1.5)
tf = box1.text_frame; tf.word_wrap = True; tf.margin_left=Inches(0.25); tf.margin_top=Inches(0.2); tf.margin_right=Inches(0.25)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Logistic Regression — baseline / scorecard"
r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = BLUE; r.font.name = FONT
for line in [
    "Features: 178 numeric + 16 categorical (194 total) — full engineered feature set",
    "Preprocessing: median imputation + standardization (numeric); most-frequent imputation + one-hot encoding (categorical)",
    "Imbalance handling: class_weight = 'balanced'",
    "Tuning: GridSearchCV, 3-fold, scoring = ROC-AUC, over C ∈ {0.01 … 3}",
    "Best C = 0.01",
    "Why keep it: fully transparent coefficients — usable as a regulator-facing scorecard",
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    r = p.add_run(); r.text = "▪ " + line; r.font.size = Pt(13); r.font.color.rgb = INK; r.font.name = FONT

# LightGBM card
box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.55), col_w, Inches(5.1))
box2.fill.solid(); box2.fill.fore_color.rgb = WHITE; box2.line.color.rgb = ORANGE; box2.line.width = Pt(1.5)
tf = box2.text_frame; tf.word_wrap = True; tf.margin_left=Inches(0.25); tf.margin_top=Inches(0.2); tf.margin_right=Inches(0.25)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "LightGBM — recommended production model"
r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = ORANGE; r.font.name = FONT
for line in [
    "Features: same 194-feature set — categoricals used natively, no manual encoding",
    "Preprocessing: none required — native missing-value handling",
    "Imbalance handling: scale_pos_weight = 11.4 (neg/pos ratio)",
    "Tuning: RandomizedSearchCV, 20 candidates x 3-fold, scoring = ROC-AUC",
    "Tuned: num_leaves, learning_rate, min_child_samples, subsample, colsample_bytree, reg_alpha/lambda",
    "Why: captures non-linear interactions (e.g. ORGANIZATION_TYPE x income ratios) a linear model can't",
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    r = p.add_run(); r.text = "▪ " + line; r.font.size = Pt(13); r.font.color.rgb = INK; r.font.name = FONT
footer(s, 6)

# ============================================================ SLIDE 7 — PERFORMANCE
s = add_slide(); set_bg(s)
header(s, "06 — MODEL PERFORMANCE", "LightGBM edges out the baseline on both credit-scoring metrics")
picture_fit(s, FIG / "roc_curve_comparison.png", Inches(0.55), Inches(1.55), Inches(5.6), Inches(4.7))

tx = Inches(6.5)
stat_tile(s, tx, Inches(1.6), Inches(2.9), Inches(1.3), f"{lr['auc']:.3f}", "Logistic Regression — Test AUC", BLUE)
stat_tile(s, tx + Inches(3.1), Inches(1.6), Inches(2.9), Inches(1.3), f"{lgb['auc']:.3f}", "LightGBM — Test AUC", ORANGE)
stat_tile(s, tx, Inches(3.1), Inches(2.9), Inches(1.3), f"{lr['ks']:.3f}", "Logistic Regression — KS", BLUE)
stat_tile(s, tx + Inches(3.1), Inches(3.1), Inches(2.9), Inches(1.3), f"{lgb['ks']:.3f}", "LightGBM — KS", ORANGE)

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, Inches(4.7), Inches(6.0), Inches(1.9))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.15); tf.margin_right=Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = ("Both models lean on the same strong engineered features (EXT_SOURCE, CREDIT_TERM, "
                            "bureau/previous-application aggregates); LightGBM's edge comes from non-linear "
                            "feature interactions a linear model cannot represent. Recommendation: LightGBM "
                            "for production scoring, Logistic Regression retained as an interpretable "
                            "challenger / regulatory fallback.")
r.font.size = Pt(13); r.font.color.rgb = INK; r.font.name = FONT
footer(s, 7)

# ============================================================ SLIDE 8 — BUSINESS IMPACT
s = add_slide(); set_bg(s)
header(s, "07 — BUSINESS IMPACT", "Catch most bad debt by declining only the riskiest few %")
picture_fit(s, FIG / "gains_curve.png", Inches(0.55), Inches(1.55), Inches(6.0), Inches(4.9))
tb = textbox(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.3), "Decile gains analysis (LightGBM score)", size=14, bold=True, color=RED)
bullets(s, Inches(6.9), Inches(2.0), Inches(5.9), Inches(3.0), [
    ("Decline riskiest 10%:", " captures 37.7% of all future defaults, affects only 7.6% of good payers"),
    ("Decline riskiest 20%:", " captures 57.1% of defaults, affects 16.7% of good payers"),
    ("Beyond ~decile 3-4:", " the trade-off worsens fast — remaining pool is genuinely low-risk"),
], size=14, space_after=14)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.9), Inches(5.9), Inches(1.9))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.15); tf.margin_right=Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Recommendation: "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = AQUA; r.font.name = FONT
r2 = p.add_run()
r2.text = ("decline (or route to a smaller/secured starter product) only the top 1-2 riskiest score deciles, "
           "where the ratio of defaults-avoided to good-payers-affected is most favorable (~5:1 and ~3.4:1). "
           "This is how the model directly serves the brief's goal: cut bad debt without mass-rejecting "
           "creditworthy customers.")
r2.font.size = Pt(13); r2.font.color.rgb = INK; r2.font.name = FONT
footer(s, 8)

# ============================================================ SLIDE 9 — RECOMMENDATION POLICY
s = add_slide(); set_bg(s)
header(s, "08 — BUSINESS RECOMMENDATION", "Translate the score into risk-tiered loan terms, not one cutoff")

rows = [
    ["Risk Tier", "Score decile", "Default rate", "Principal ceiling", "Max tenor", "Down payment", "Repayment calendar"],
    ["1 — Low", "9–10", "~1–2%", "Up to 5x income", "36–48 mo", "None", "Standard monthly, fast-track"],
    ["2 — Medium", "5–8", "~3–7%", "Up to 4x income", "24–36 mo", "10–20%", "Standard monthly"],
    ["3 — Elevated", "3–4", "~7–10%", "Up to 3x income", "12–24 mo", "20–30%", "Bi-weekly instalments"],
    ["4 — High", "1–2", "> 15%", "Up to 2x income (starter/secured)", "6–12 mo", "30%+ / guarantor", "Weekly; graduate after clean history"],
]
n_rows, n_cols = len(rows), len(rows[0])
table_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.55), Inches(1.6), Inches(12.2), Inches(3.0))
table = table_shape.table
col_widths = [1.4, 1.1, 1.1, 2.5, 1.1, 1.5, 3.5]
for i, w in enumerate(col_widths):
    table.columns[i].width = Inches(w)
tier_colors = [RED, ORANGE, RGBColor(0xED, 0xA1, 0x00), AQUA]
for r_idx, row in enumerate(rows):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        for run in para.runs:
            run.font.size = Pt(11.5)
            run.font.name = FONT
            if r_idx == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
            else:
                run.font.color.rgb = INK
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

textbox(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.8),
        "Fairness note: CODE_GENDER shows a real gap in the data (Male 10.1% vs Female 7.0% default) but is "
        "excluded as a scoring/pricing lever — the underlying risk is already captured through income "
        "stability, external bureau score, and employment type. A correlate is not automatically a lever.",
        size=12.5, color=INK_SECONDARY, italic=True)
textbox(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.0),
        "Next steps: validate on an out-of-time sample (not just a random holdout), add reject-inference for "
        "the declined population (models are trained only on historically-approved customers), and A/B test "
        "the repayment-calendar policy against the current flat policy before full rollout.",
        size=12.5, color=INK_SECONDARY)
footer(s, 9)

# ============================================================ SLIDE 10 — REFERENCES
s = add_slide(); set_bg(s)
header(s, "09 — REFERENCES & REPOSITORY", "Full analysis, code, and reproducible notebook")

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.3))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "GitHub repository: "; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = RED; r.font.name = FONT
r2 = p.add_run(); r2.text = "github.com/glecyhaloho/home-credit-default-risk-scoring"
r2.font.size = Pt(16); r2.font.color.rgb = INK; r2.font.name = FONT
p2 = tf.add_paragraph()
r3 = p2.add_run(); r3.text = "Contains the full executed notebook (.ipynb), feature-engineering & training source code, and saved model artifacts."
r3.font.size = Pt(13); r3.font.color.rgb = INK_SECONDARY; r3.font.name = FONT

textbox(s, Inches(0.55), Inches(3.3), Inches(12.2), Inches(0.4), "References", size=15, bold=True, color=INK)
bullets(s, Inches(0.55), Inches(3.75), Inches(12.2), Inches(2.6), [
    "Kaggle — \"Home Credit Default Risk\" competition dataset & problem description (source of application_{train,test}, bureau, previous_application, POS_CASH_balance, credit_card_balance, installments_payments tables)",
    "Rakamin Academy x Home Credit Indonesia — Virtual Internship Task 5 brief: Dataset Description & PPT guidelines",
    "scikit-learn documentation — Logistic Regression, GridSearchCV, ColumnTransformer",
    "LightGBM documentation — gradient boosting framework, categorical feature handling",
    "Standard credit-scoring practice — KS statistic, decile/gains-table analysis for accept/reject policy design",
], size=13.5, space_after=10)
footer(s, 10)

prs.save(ROOT / "reports" / "Home_Credit_Task5_Presentation.pptx")
print("Saved PPT:", ROOT / "reports" / "Home_Credit_Task5_Presentation.pptx")
print("Total slides:", len(prs.slides.__iter__.__self__._sldIdLst))
